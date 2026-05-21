"""Generate Lab_Document_Query.pptx — short non-technical overview."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x11, 0x18, 0x27)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xEE, 0xF2, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=DARK, align=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        if align is not None:
            p.alignment = align
    return box


def add_bullets(slide, left, top, width, height, items, *, size=16, color=DARK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, sub) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(8)
        r = p.add_run()
        r.text = "•  "
        r.font.size = Pt(size)
        r.font.color.rgb = BLUE
        r.font.bold = True
        r2 = p.add_run()
        r2.text = head
        r2.font.size = Pt(size)
        r2.font.bold = True
        r2.font.color.rgb = color
        if sub:
            r3 = p.add_run()
            r3.text = "  —  " + sub
            r3.font.size = Pt(size - 2)
            r3.font.color.rgb = GRAY
    return box


def add_accent_bar(slide, left, top, width, height, color=BLUE):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide 1 — Title ──────────────────────────────────────────────────────────
blank = prs.slide_layouts[6]
s1 = prs.slides.add_slide(blank)
set_bg(s1, LIGHT_BG)

add_accent_bar(s1, Inches(0), Inches(0), Inches(13.333), Inches(0.4))

add_text(s1, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
         "Lab Document Query", size=54, bold=True, color=DARK)
add_text(s1, Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.8),
         "AI-drevet innsikt fra rapporter og forskning om barn og unge",
         size=24, color=BLUE)
add_text(s1, Inches(0.8), Inches(4.3), Inches(11.7), Inches(1.5),
         "Gjør hundrevis av sider med rapporter, artikler og nettressurser\n"
         "om barn og unges hverdag til søkbar, sporbar innsikt — på sekunder.",
         size=18, color=GRAY)

# ── Slide 2 — Hva brukeren kan gjøre ─────────────────────────────────────────
s2 = prs.slides.add_slide(blank)
set_bg(s2, WHITE)
add_accent_bar(s2, Inches(0), Inches(0), Inches(13.333), Inches(0.4))

add_text(s2, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8),
         "Hva kan brukeren gjøre?", size=32, bold=True, color=DARK)

# Two columns
col_w = Inches(5.8)
col_top = Inches(1.9)
col_h = Inches(5.0)

# Column 1 — Spør direkte
add_accent_bar(s2, Inches(0.8), col_top, Inches(0.12), col_h, color=BLUE)
add_text(s2, Inches(1.05), col_top, col_w, Inches(0.6),
         "Spør direkte", size=22, bold=True, color=BLUE)
add_text(s2, Inches(1.05), Inches(2.45), col_w, Inches(0.7),
         "Still et spørsmål om et konkret tema og få et\n"
         "svar med kildehenvisninger til rapport og side.",
         size=15, color=DARK)
add_bullets(s2, Inches(1.05), Inches(3.6), col_w, Inches(3.2), [
    ("Eksempel:", "«Hva sier rapportene om skolefravær?»"),
    ("Kilder", "Hver påstand peker tilbake til original tekst"),
    ("Filter", "Avgrens på rapport, målgruppe, segment, år"),
], size=14)

# Column 2 — Aggregert analyse
add_accent_bar(s2, Inches(7.0), col_top, Inches(0.12), col_h, color=RGBColor(0x5B, 0x21, 0xB6))
add_text(s2, Inches(7.25), col_top, col_w, Inches(0.6),
         "Aggregert analyse", size=22, bold=True, color=RGBColor(0x5B, 0x21, 0xB6))
add_text(s2, Inches(7.25), Inches(2.45), col_w, Inches(0.7),
         "Syntetiser funn på tvers av hele biblioteket\n"
         "— én knapp, ferdig rapport.",
         size=15, color=DARK)
add_bullets(s2, Inches(7.25), Inches(3.6), col_w, Inches(3.2), [
    ("Problemer", "Hvilke utfordringer sliter unge med?"),
    ("Kritiske øyeblikk", "Vendepunkter i unges liv"),
    ("Personas", "Syntetiserte brukerprofiler"),
    ("Fri analyse", "Åpne spørsmål på tvers av alle dokumenter"),
], size=14)

# ── Slide 3 — Verdi ──────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(blank)
set_bg(s3, LIGHT_BG)
add_accent_bar(s3, Inches(0), Inches(0), Inches(13.333), Inches(0.4))

add_text(s3, Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8),
         "Hva får organisasjonen igjen?", size=32, bold=True, color=DARK)

add_bullets(s3, Inches(1.0), Inches(2.0), Inches(11.5), Inches(5.0), [
    ("Raskt overblikk",
     "Få ut hovedfunnene fra 100+ sider på minutter, ikke dager"),
    ("Sporbarhet",
     "Hver påstand kobles til rapport og sidetall — ingen gjetting"),
    ("Delbart resultat",
     "Last ned analysen som ferdig .docx-rapport"),
    ("Fleksibelt",
     "Bytt mellom kunnskapsbibliotek — f.eks. forskningsrapporter eller nettinnhold"),
    ("Skalerbart",
     "Nye kilder (PDF, artikler, nettsider) kan legges inn etter hvert"),
], size=18)

add_text(s3, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5),
         "Målgruppe: produktteam, designere og beslutningstakere som trenger ungdomsinnsikt.",
         size=14, color=GRAY)

out = "Lab_Document_Query.pptx"
prs.save(out)
print(f"Wrote {out}")
